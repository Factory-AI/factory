"""
generate_slides.py
Run: python generate_slides.py
Output: jeval.pptx  — open directly in Google Slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK  = RGBColor(0, 0, 0)
WHITE  = RGBColor(255, 255, 255)
GREY   = RGBColor(160, 160, 160)
LGREY  = RGBColor(230, 230, 230)
MONO   = "Courier New"
SANS   = "Courier New"

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def slide():
    return prs.slides.add_slide(BLANK)


def box(sl, left, top, width, height,
        text="", size=14, bold=False, color=BLACK,
        font=SANS, align=PP_ALIGN.LEFT, wrap=True):
    txBox = sl.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def label(sl, text):
    box(sl, 0.8, 0.5, 11, 0.4, text, size=9, color=GREY, font=MONO)


def hline(sl, top):
    line = sl.shapes.add_connector(1, Inches(0.8), Inches(top), Inches(12.5), Inches(top))
    line.line.color.rgb = LGREY
    line.line.width = Pt(0.5)


def codebox(sl, left, top, width, height, text):
    shape = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 248, 248)
    shape.line.color.rgb = LGREY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = MONO
    run.font.size = Pt(11)
    run.font.color.rgb = BLACK


def multiline(sl, left, top, width, lines, size=13, gap=0.38, color=BLACK, bold_first=False):
    for i, line in enumerate(lines):
        b = bold_first and i == 0
        box(sl, left, top + i * gap, width, gap + 0.1,
            line, size=size, bold=b, color=color if not (b) else BLACK)


# ── Slide 1: Title ────────────────────────────────────────────
s = slide()
box(s, 0.8, 1.2, 10, 0.5, "FACTORY AI · MARCH 2026", size=9, color=GREY)
box(s, 0.8, 1.9, 10, 1.2, "jeval", size=52, bold=True)
box(s, 0.8, 3.3, 8,  0.6,
    "JEPA-based semantic fidelity verification for Droid memory compression",
    size=18)
box(s, 0.8, 5.8, 10, 0.4,
    "github.com/Pshyam17/factory  ·  feat/jeval-memory-compression",
    size=9, color=GREY)

# ── Slide 2: Problem ─────────────────────────────────────────
s = slide()
label(s, "THE PROBLEM")
box(s, 0.8, 1.1, 11, 1.2,
    "Artifact tracking scores 2.45/5 across every compression method.",
    size=26, bold=True)

rows = [("Factory (best)", 2.45), ("Anthropic", 2.33), ("OpenAI", 2.19)]
for i, (name, score) in enumerate(rows):
    y = 2.8 + i * 0.6
    box(s, 0.8, y, 2.0, 0.4, name, size=13, color=GREY)
    bar = s.shapes.add_shape(1, Inches(2.9), Inches(y + 0.1),
                              Inches(score / 5 * 7), Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()
    box(s, 10.1, y, 0.8, 0.4, str(score), size=13, bold=True)

hline(s, 4.8)
box(s, 0.8, 4.9, 11, 0.8,
    '"Artifact tracking may need dedicated state tracking beyond summarization."\n— Factory AI, Dec 2025',
    size=11, color=GREY)

# ── Slide 3: Core Idea ───────────────────────────────────────
s = slide()
label(s, "THE CORE IDEA")
box(s, 0.8, 1.1, 10, 1.0,
    "Use JEPA's prediction residual as a pre-hoc fidelity oracle.",
    size=24, bold=True)
codebox(s, 0.8, 2.4, 11.5, 1.8,
    "# standard JEPA — training signal\n"
    "L = ||predictor(enc(compressed)) - enc(original)||²\n\n"
    "# jeval — fidelity measurement\n"
    "EPE(T, C) = sum((pred(enc(C)) - enc(T))²) / 4")
multiline(s, 0.8, 4.6, 11,
    ["· No prior work uses the JEPA residual as an external fidelity metric",
     "· EPE = 0 → perfect reconstruction   EPE = 1 → total semantic loss",
     "· Sensitive to role reversal, negation, causal inversion — cosine is not"],
    size=13, color=GREY)

# ── Slide 4: Architecture ────────────────────────────────────
s = slide()
label(s, "ARCHITECTURE")
steps = [
    ("01", "SEGMENT",  "bullet points + headers → individual entries"),
    ("02", "CLASSIFY", "zero-shot NLI → FACTUAL / CAUSAL / ENTITY / TEMPORAL / CONTRASTIVE / BACKGROUND"),
    ("03", "EPE",      "trained JEPA predictor → embedding prediction error per segment"),
    ("04", "BUDGET",   "z-score + content type → compression tier per segment"),
    ("05", "COMPRESS", "Mistral via NVIDIA NIM · budget=1.0 verbatim · budget=0.3 aggressive"),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.1 + i * 1.0
    box(s, 0.8, y, 0.5, 0.5, num,  size=11, color=GREY)
    box(s, 1.4, y, 1.6, 0.5, title, size=13, bold=True)
    box(s, 3.2, y, 9.5, 0.5, desc,  size=13, color=GREY)
    if i < len(steps) - 1:
        hline(s, y + 0.75)

# ── Slide 5: Design Decisions ────────────────────────────────
s = slide()
label(s, "KEY DESIGN DECISIONS")
decisions = [
    ("Why freeze the encoder?",
     "A moving target makes EPE uncalibrated. Freezing all-mpnet-base-v2 gives a fixed semantic geometry — EPE has one meaning across all sessions."),
    ("Why z-scores not raw thresholds?",
     "Raw thresholds are hardcoded to one trained predictor. Z-scores normalize against the session's own EPE distribution — self-calibrating."),
    ("Why artifact pattern detection?",
     "Low EPE ≠ low importance. File paths are predictable so the predictor assigns them low EPE — but src/auth/refresh.ts is critical. Pattern detection overrides EPE."),
    ("Why sum not mean in MSE?",
     "mean() / 768 dims = 0.003 for orthogonal vectors — indistinguishable from verbatim. sum() preserves the full signal. /4 normalizes to [0,1]."),
]
for i, (q, a) in enumerate(decisions):
    y = 1.1 + i * 1.4
    box(s, 0.8, y,       11, 0.45, q, size=13, bold=True)
    box(s, 0.8, y + 0.45, 11, 0.7, a, size=12, color=GREY)

# ── Slide 6: Training ────────────────────────────────────────
s = slide()
label(s, "TRAINING THE PREDICTOR")
box(s, 0.8, 1.1, 5, 0.4, "Setup", size=14, bold=True)
multiline(s, 0.8, 1.6, 5.5,
    ["Frozen encoder: all-mpnet-base-v2",
     "Predictor: 3-layer Pre-LN transformer",
     "5,000 synthetic pairs · 30 epochs · A100",
     "3 compression strategies:",
     "  · truncate (cut end)",
     "  · word dropout (30–50%)",
     "  · abstractify (replace specifics)"],
    size=12, color=GREY)
box(s, 6.8, 1.1, 5.5, 0.4, "Sanity check after training", size=14, bold=True)
codebox(s, 6.8, 1.6, 5.7, 2.2,
    "verbatim    EPE = 0.0060  ✓\n"
    "abstractify EPE = 0.1023  ✓\n"
    "abstractify EPE = 0.1687  ✓\n\n"
    "17–28× separation")

# ── Slide 7: Results ─────────────────────────────────────────
s = slide()
label(s, "RESULTS")
box(s, 0.8, 1.0, 11, 0.4, "Artifact survival — 3 iterative compression rounds", size=13, bold=True)
rows2 = [("Stage 1","84 tokens","6/10","future artifacts not yet written"),
         ("Stage 2","175 tokens","9/10",""),
         ("Stage 3","261 tokens","10/10","all critical artifacts preserved ✓")]
for i,(st,tk,sc,n) in enumerate(rows2):
    y = 1.5 + i*0.5
    box(s, 0.8, y, 1.0, 0.4, st, size=12, color=GREY)
    box(s, 2.0, y, 1.5, 0.4, tk, size=12, color=GREY)
    box(s, 3.7, y, 1.2, 0.4, sc, size=12, bold=(sc=="10/10"))
    box(s, 5.1, y, 7.0, 0.4, n,  size=11, color=GREY)
    hline(s, y + 0.45)

box(s, 0.8, 3.4, 11, 0.4, "Artifact tracking vs Factory baseline", size=13, bold=True)
results = [("jeval","4.75",True),("Factory","2.45",False),("Anthropic","2.33",False),("OpenAI","2.19",False)]
for i,(m,a,bold) in enumerate(results):
    y = 3.9 + i*0.6
    box(s, 0.8, y, 1.5, 0.45, m, size=13, bold=bold)
    bar_w = float(a) / 5 * 7
    bar = s.shapes.add_shape(1, Inches(2.5), Inches(y+0.12),
                              Inches(bar_w), Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK if bold else LGREY
    bar.line.fill.background()
    box(s, 9.8, y, 0.8, 0.45, a, size=13, bold=bold)
    if bold:
        box(s, 10.7, y, 1.5, 0.45, "+94%", size=9, color=GREY)

# ── Slide 8: Next Steps ──────────────────────────────────────
s = slide()
label(s, "NEXT STEPS")
nexts = [
    ("01","Real session validation",
     "Retrain predictor on Droid production sessions. Run probe eval across 20-30 sessions with confidence intervals."),
    ("02","Ablation study",
     "Quantify contribution of each component: EPE alone, z-scores alone, artifact patterns alone, full system."),
    ("03","Incremental architecture",
     "Classify + compute EPE per entry as it's written. PreCompact reads cached values. Hook latency → <500ms."),
    ("04","EMNLP 2026 System Demos",
     "Deadline July 4, 2026. Working system + empirical results on published benchmark."),
]
for i,(num,title,desc) in enumerate(nexts):
    y = 1.1 + i * 1.4
    box(s, 0.8, y,        0.5, 0.45, num,   size=11, color=GREY)
    box(s, 1.4, y,        3.2, 0.45, title, size=13, bold=True)
    box(s, 1.4, y + 0.5,  11,  0.7,  desc,  size=12, color=GREY)

# ── Slide 9: Install ─────────────────────────────────────────
s = slide()
label(s, "INSTALL")
codebox(s, 0.8, 1.2, 11.5, 3.2,
    "# install\n"
    "pip install -e examples/jeval-memory-compression\n\n"
    "# set key\n"
    "export NVIDIA_API_KEY=nvapi-...\n\n"
    "# register hook in ~/.factory/settings.json\n"
    '"PreCompact": [{ "type": "command",\n'
    '  "command": "python3 $PROJECT_DIR/...precompact_jeval.py" }]')
multiline(s, 0.8, 5.0, 12,
    ["github.com/Pshyam17/factory   ·   feat/jeval-memory-compression   ·   PR #785 open"],
    size=10, color=GREY)

# ── Save ─────────────────────────────────────────────────────
prs.save("jeval.pptx")
print("✓ saved jeval.pptx")
print("  upload to Google Slides via File → Import slides")